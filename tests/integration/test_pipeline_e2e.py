"""End-to-end: the full graph on repair.pdf with all agents faked.

Exercises ingest -> detect -> Send fan-out -> gather -> content -> translate
-> assemble -> qa against the REAL PDF, producing a real .pptx that must
re-open. No network: vision verifier, content agent, critic, and translator
are all fakes. This is the reducer/orchestration regression guard.
"""

import os

import pytest
from pptx import Presentation

from pdfdeck.agents.content_agent import (
    CritiqueReport,
    DraftBundle,
    DraftedSlide,
    FakeContentAgent,
    FakeFidelityCritic,
    Outline,
    PlannedSlide,
)
from pdfdeck.agents.vision_verifier import FakeVisionVerifier, Verdict, VerificationResult
from pdfdeck.pipeline import convert_pdf

pytestmark = pytest.mark.integration

FIXTURES = os.path.join(os.path.dirname(__file__), "..", "fixtures")
PDF = os.path.join(FIXTURES, "repair.pdf")


class _AcceptAllVerifier:
    """Accepts every crop; unbounded (Send fan-out reuses one verifier)."""

    def __init__(self):
        self.calls = 0

    def verify(self, image_path, expected_caption, kind):
        self.calls += 1
        return VerificationResult(verdict=Verdict.ACCEPT, reason="ok")


def _content_agent():
    outline = Outline(topic="Tissue Repair", slides=[
        PlannedSlide(title="Tissue Repair", kind="title"),
        PlannedSlide(title="Granulation Tissue", kind="content", span_ids=["p0_b0"]),
    ])
    draft = DraftBundle(slides=[
        DraftedSlide(title="Granulation Tissue",
                     bullets=["New capillaries and fibroblasts appear by days 3-5."],
                     source_span_ids=["p0_b0"], kind="content"),
    ])
    return FakeContentAgent(outline, [draft])


class _FakeTranslator:
    def is_configured(self):
        return True

    def translate_batch(self, texts, target, source="en"):
        return [f"[{target}] {t}" for t in texts]


def test_full_pipeline_produces_openable_deck(tmp_path):
    result = convert_pdf(
        PDF,
        target_language=None,
        vision_enabled=True,
        run_dir=str(tmp_path),
        verifier=_AcceptAllVerifier(),
        content_agent=_content_agent(),
        critic=FakeFidelityCritic([CritiqueReport(approved=True)]),
        translator=_FakeTranslator(),
    )

    # Deck exists and re-opens.
    assert os.path.exists(result.output_path)
    prs = Presentation(result.output_path)
    slides = list(prs.slides)
    assert len(slides) >= 3  # title + content + figures

    # All 8 calibrated figures made it through the fan-out (conservation).
    assert len(result.figures) == 8

    # Every figure slide carries a picture shape (aspect-preserved, not text).
    pic_slides = 0
    for slide in slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                pic_slides += 1
                break
    assert pic_slides >= 8

    # QA report was written.
    assert os.path.exists(os.path.join(str(tmp_path), "qa_report.json"))


def test_translation_applied(tmp_path):
    result = convert_pdf(
        PDF,
        target_language="tr",
        vision_enabled=False,  # deterministic figures, faster
        run_dir=str(tmp_path),
        verifier=_AcceptAllVerifier(),
        content_agent=_content_agent(),
        critic=FakeFidelityCritic([CritiqueReport(approved=True)]),
        translator=_FakeTranslator(),
    )
    text_slides = [s for s in result.slides if s.slide_type == "text"]
    assert text_slides
    assert all(b.startswith("[tr] ") for s in text_slides for b in s.bullets)


def test_extra_languages_emits_second_deck_from_one_run(tmp_path):
    """English primary + extra_languages=['tr'] yields two openable decks from a
    single (expensive) run: the English content is translated for the extra."""
    result = convert_pdf(
        PDF,
        target_language=None,  # English primary
        vision_enabled=False,  # deterministic figures, faster
        run_dir=str(tmp_path),
        extra_languages=["tr", "en"],  # 'en' is a no-op and must be skipped
        verifier=_AcceptAllVerifier(),
        content_agent=_content_agent(),
        critic=FakeFidelityCritic([CritiqueReport(approved=True)]),
        translator=_FakeTranslator(),
    )

    # English deck (primary) is untranslated.
    text_slides = [s for s in result.slides if s.slide_type == "text"]
    assert text_slides
    assert not any(b.startswith("[tr] ") for s in text_slides for b in s.bullets)

    # Exactly one extra deck (tr); 'en' was skipped.
    assert list(result.extra_outputs) == ["tr"]
    tr_path = result.extra_outputs["tr"]
    assert os.path.exists(tr_path)
    assert tr_path != result.output_path
    assert os.path.exists(result.output_path)

    # The extra deck re-opens with the same slide count as the English one.
    assert len(list(Presentation(tr_path).slides)) == len(list(Presentation(result.output_path).slides))


def test_no_vision_mode_makes_no_verifier_calls(tmp_path):
    v = _AcceptAllVerifier()
    convert_pdf(
        PDF,
        vision_enabled=False,
        run_dir=str(tmp_path),
        verifier=v,
        content_agent=_content_agent(),
        critic=FakeFidelityCritic([CritiqueReport(approved=True)]),
        translator=_FakeTranslator(),
    )
    assert v.calls == 0
