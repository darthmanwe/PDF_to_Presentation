from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import io
from typing import List, Dict, Any, Optional
from PIL import Image

class PresentationBuilder:
    """Handles creation of PowerPoint presentations from detailed medical content"""
    
    def __init__(self):
        self.presentation = Presentation()
        self._setup_default_layouts()
    
    def _setup_default_layouts(self):
        """Setup default slide layouts"""
        # Use the default slide layout (title and content)
        self.slide_layout = self.presentation.slide_layouts[1]  # Title and Content layout
    
    def create_presentation(self, slides_data: List[Dict[str, Any]], 
                          output_path: str, 
                          title: str = "Medical Presentation",
                          subtitle: str = "Generated from PDF") -> str:
        """Create a complete presentation from slides data"""
        
        try:
            # Create title slide
            self._create_title_slide(title, subtitle)
            
            # Create content slides
            for slide_data in slides_data:
                self._create_content_slide(slide_data)
            
            # Save the presentation
            self.presentation.save(output_path)
            return output_path
            
        except Exception as e:
            raise Exception(f"Error creating presentation: {e}")
    
    def _create_title_slide(self, title: str, subtitle: str):
        """Create the title slide with medical topic only (NO presenter info)"""
        # Use the title slide layout
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        
        # Set title - use the actual medical topic
        title_shape = slide.shapes.title
        if not title or title.strip() == "" or "untitled" in title.lower() or "tmp" in title.lower():
            # Use a proper medical title if the provided title is empty or generic
            title = "Medical Education Presentation"
            print("🔍 DEBUG: Using default medical title for empty/generic title")
        
        title_shape.text = title
        self._format_title(title_shape)
        print(f"🔍 DEBUG: Title slide created with title: '{title}'")
        
        # Set subtitle - medical education focus
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        self._format_subtitle(subtitle_shape)
    
    def _create_text_slide(self, slide, slide_data: Dict[str, Any]):
        """Create a text-only slide with paragraphs (NO bullet points)"""
        content_shape = slide.placeholders[1]
        content = slide_data.get('content', [])
        print(f"🔍 DEBUG: Raw content from slide data: {content}")
        print(f"🔍 DEBUG: Content type: {type(content)}")
        print(f"🔍 DEBUG: Content length: {len(content) if isinstance(content, list) else 'N/A'}")
        
        # Process content as paragraphs (NO bullet points)
        if isinstance(content, list):
            # Content is a list of paragraphs
            processed_content = []
            for item in content:
                if item and str(item).strip():
                    processed_item = str(item).strip()
                    processed_content.append(processed_item)
                    print(f"🔍 DEBUG: Added content paragraph: '{processed_item[:100]}...'")
            
            # Add paragraphs to PowerPoint (NO bullet points)
            if processed_content:
                self._add_paragraphs_to_shape(content_shape, processed_content)
                print(f"🔍 DEBUG: Added {len(processed_content)} paragraphs to slide")
            else:
                content_shape.text = "No content available"
                print("🔍 DEBUG: No content items found, using fallback")
        else:
            # Handle string content
            content_str = str(content)
            if content_str.strip():
                content_shape.text = content_str.strip()
                print(f"🔍 DEBUG: String content set: {content_str.strip()[:100]}...")
            else:
                content_shape.text = "No content available"
                print("🔍 DEBUG: Empty string content, using fallback")
    
    def _create_image_slide(self, slide, slide_data: Dict[str, Any], pdf_content: Dict[str, Any] = None):
        """Create an image-focused slide with maximum 2 images from the same page"""
        # Add images if present and PDF content is available
        if pdf_content and 'page_images' in slide_data and 'page_number' in slide_data:
            page_number = slide_data['page_number']
            page_images = slide_data['page_images']
            print(f"Adding {len(page_images)} images from page {page_number}")
            
            # Calculate positions for maximum 2 images
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            
            # Position images side by side or stacked
            if len(page_images) == 1:
                # Single image - center it
                image_width = Inches(4)
                image_height = Inches(3)
                left = Inches(3)  # Center horizontally
                top = Inches(2.5)  # Center vertically
                
                self._add_single_image(slide, page_images[0], left, top, image_width, image_height, page_number, 0)
                
            elif len(page_images) == 2:
                # Two images - side by side
                image_width = Inches(3.5)
                image_height = Inches(2.5)
                
                # Left image
                left1 = Inches(1)
                top1 = Inches(2.5)
                self._add_single_image(slide, page_images[0], left1, top1, image_width, image_height, page_number, 0)
                
                # Right image
                left2 = Inches(5.5)
                top2 = Inches(2.5)
                self._add_single_image(slide, page_images[1], left2, top2, image_width, image_height, page_number, 1)
        
        # Add image descriptions if present
        if 'image_descriptions' in slide_data and slide_data['image_descriptions']:
            descriptions = slide_data['image_descriptions']
            print(f"Adding {len(descriptions)} image descriptions")
            
            # Add descriptions below images
            for i, description in enumerate(descriptions):
                try:
                    if len(descriptions) == 1:
                        # Single image - center description
                        left = Inches(2)
                        top = Inches(5.5)
                        width = Inches(6)
                    else:
                        # Two images - position below corresponding image
                        left = Inches(1 + i * 4.5)  # 1 inch margin + 4.5 inch spacing
                        top = Inches(5.5)
                        width = Inches(3.5)
                    
                    height = Inches(0.5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    p = text_frame.paragraphs[0]
                    p.text = description
                    p.font.size = Pt(10)
                    p.font.italic = True
                    p.font.color.rgb = RGBColor(89, 89, 89)
                    p.alignment = PP_ALIGN.CENTER
                    
                except Exception as e:
                    print(f"Error adding image description {i+1}: {e}")
    
    def _add_single_image(self, slide, img_data, left, top, width, height, page_number, img_index):
        """Helper method to add a single image to slide"""
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(img_data['data']))
            
            # Save temporarily
            temp_path = f"temp_image_{page_number}_{img_index}.png"
            image.save(temp_path)
            
            # Add to slide
            slide.shapes.add_picture(temp_path, left, top, width, height)
            
            # Clean up
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
            # Minimal logging - just success message
            print(f"✓ Image {img_index+1} from page {page_number} added")
            
        except Exception as e:
            print(f"✗ Error adding image {img_index+1} from page {page_number}: {e}")
    
    def _create_mixed_slide(self, slide, slide_data: Dict[str, Any]):
        """Create a mixed slide with both text and image"""
        # Add text content first
        self._create_text_slide(slide, slide_data)
        
        # Add image if present
        if 'image_path' in slide_data and slide_data['image_path']:
            print(f"Adding image to mixed slide: {slide_data['image_path']}")
            self._add_image_to_slide(slide, slide_data['image_path'])
    
    def _format_title(self, title_shape):
        """Format the main title slide text - larger size"""
        title_frame = title_shape.text_frame
        title_frame.clear()
        
        p = title_frame.paragraphs[0]
        p.text = title_shape.text
        p.font.size = Pt(44)  # Large size for main title
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
    
    def _format_slide_title(self, title_shape):
        """Format slide titles - smaller size"""
        title_frame = title_shape.text_frame
        title_frame.clear()
        
        p = title_frame.paragraphs[0]
        p.text = title_shape.text
        p.font.size = Pt(28)  # Smaller size for slide titles
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT  # Left align for content slides
    
    def _format_subtitle(self, subtitle_shape):
        """Format the subtitle text"""
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.clear()
        
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle_shape.text
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = RGBColor(89, 89, 89)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_detailed_bullet_points_to_shape(self, content_shape, bullet_points: List[str]):
        """Add detailed bullet points directly to the content shape"""
        print(f"Adding {len(bullet_points)} detailed bullet points to shape")
        
        # Clear the text frame
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add each detailed bullet point as a separate paragraph
        for i, point in enumerate(bullet_points):
            if i == 0:
                # Use the first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraph for subsequent points
                p = text_frame.add_paragraph()
            
            # Set the text and bullet level
            p.text = point
            p.level = 0  # Main bullet level
            
            # Set font properties for detailed content
            p.font.size = Pt(16)  # Smaller font for detailed content
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            print(f"Added detailed bullet point {i+1}: '{point}'")
        
        print(f"Successfully added {len(bullet_points)} detailed bullet points to text frame")
    
    def _add_paragraphs_to_shape(self, content_shape, paragraphs: List[str]):
        """Add paragraphs directly to the content shape (NO bullet points)"""
        print(f"Adding {len(paragraphs)} paragraphs to shape")
        
        # Clear the text frame
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add each paragraph as a separate paragraph (NO bullet points)
        for i, paragraph in enumerate(paragraphs):
            if i == 0:
                # Use the first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraph for subsequent paragraphs
                p = text_frame.add_paragraph()
            
            # Set the text (NO bullet points)
            p.text = paragraph
            
            # Set font properties for paragraphs
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            print(f"Added paragraph {i+1}: '{paragraph}'")
        
        print(f"Successfully added {len(paragraphs)} paragraphs to text frame")
    
    def _add_image_to_slide(self, slide, image_path: str):
        """Add an image to the slide"""
        try:
            # Check if image file exists
            if not os.path.exists(image_path):
                return
            
            # Calculate image position and size
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            
            # Position image on the right side
            left = Inches(6)
            top = Inches(2)
            width = Inches(3)
            height = Inches(4)
            
            # Add image to slide
            slide.shapes.add_picture(image_path, left, top, width, height)
            
        except Exception as e:
            print(f"Error adding image to slide: {e}")
    
    def _add_image_description(self, slide, description: str):
        """Add image description to the slide"""
        try:
            # Add a text box for the image description
            left = Inches(6)
            top = Inches(6.5)
            width = Inches(3)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(89, 89, 89)
            p.alignment = PP_ALIGN.CENTER
            
        except Exception as e:
            print(f"Error adding image description: {e}")
    
    def add_image_from_bytes(self, slide, image_data: bytes, image_name: str = "image.png"):
        """Add an image to slide from bytes data"""
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(image_data))
            
            # Save temporarily
            temp_path = f"temp_{image_name}"
            image.save(temp_path)
            
            # Add to slide
            self._add_image_to_slide(slide, temp_path)
            
            # Clean up
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
        except Exception as e:
            print(f"Error adding image from bytes: {e}")
    
    def create_medical_presentation(self, slides_data: List[Dict[str, Any]], 
                                  output_path: str,
                                  pdf_title: str = "Medical Document",
                                  pdf_content: Dict[str, Any] = None,
                                  include_images: bool = True) -> str:
        """Create a detailed medical presentation with proper formatting"""
        
        try:
            print(f"Creating detailed medical presentation with {len(slides_data)} slides")
            print(f"PDF Title: {pdf_title}")
            print(f"PDF content available: {pdf_content is not None}")
            
            # Create title slide with medical topic only (NO presenter info)
            print(f"🔍 DEBUG: Creating title slide with title: '{pdf_title}'")
            self._create_title_slide(
                title=pdf_title,
                subtitle="Medical Education Presentation"
            )
            print(f"🔍 DEBUG: Title slide created")
            
            # Create content slides
            for i, slide_data in enumerate(slides_data):
                print(f"Processing slide {i+1}: Title='{slide_data.get('title', 'NO TITLE')}', Type='{slide_data.get('slide_type', 'NO TYPE')}'")
                
                # Add slide number to title if not present
                if 'title' not in slide_data or not slide_data['title']:
                    slide_data['title'] = f"Slide {i + 1}"
                    print(f"🔍 DEBUG: Added fallback title: '{slide_data['title']}'")
                
                self._create_content_slide(slide_data, pdf_content)
            
            # Save presentation
            self.presentation.save(output_path)
            print(f"Presentation saved with {len(self.presentation.slides)} slides")
            return output_path
            
        except Exception as e:
            print(f"Error in create_medical_presentation: {e}")
            raise Exception(f"Error creating medical presentation: {e}")
    
    def _create_content_slide(self, slide_data: Dict[str, Any], pdf_content: Dict[str, Any] = None):
        """Create a content slide from structured slide data"""
        print(f"Creating content slide: Title='{slide_data.get('title', 'NO TITLE')}', Type='{slide_data.get('slide_type', 'NO TYPE')}'")
        
        slide = self.presentation.slides.add_slide(self.slide_layout)
        
        # Set slide title - smaller size
        title_shape = slide.shapes.title
        title = slide_data.get('title', 'Untitled Slide')
        print(f"Setting slide title: '{title}'")
        title_shape.text = title
        self._format_slide_title(title_shape)  # Smaller title formatting
        
        # Handle different slide types
        slide_type = slide_data.get('slide_type', 'text')
        print(f"Slide type: {slide_type}")
        
        if slide_type == 'text':
            self._create_text_slide(slide, slide_data)
        elif slide_type == 'image':
            self._create_image_slide(slide, slide_data, pdf_content)
        elif slide_type == 'annotation':
            self._create_annotation_slide(slide, slide_data)
        elif slide_type == 'mixed':
            self._create_mixed_slide(slide, slide_data, pdf_content)
        else:
            # Default to text slide
            self._create_text_slide(slide, slide_data)
        
        print(f"Slide creation completed for: {title}")
    
    def _create_annotation_slide(self, slide, slide_data: Dict[str, Any]):
        """Create an annotation slide with figure/table descriptions (text-only mode)"""
        content_shape = slide.placeholders[1]
        
        # Get annotation descriptions
        annotations = slide_data.get('annotation_descriptions', [])
        print(f"Adding {len(annotations)} annotation descriptions")
        
        if annotations:
            # Add each annotation as a paragraph
            for i, annotation in enumerate(annotations):
                try:
                    # Create a text box for each annotation
                    left = Inches(1)
                    top = Inches(2 + i * 1.5)  # Stack annotations vertically
                    width = Inches(8)
                    height = Inches(1.2)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    # Add the annotation text
                    p = text_frame.paragraphs[0]
                    p.text = annotation
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT
                    
                    # Add some spacing
                    p.space_after = Pt(6)
                    
                except Exception as e:
                    print(f"Error adding annotation {i+1}: {e}")
        else:
            # Fallback if no annotations
            content_shape.text = "No figure or table annotations found"
            print("No annotations found, using fallback text")
    
    def _create_image_slide(self, slide, slide_data: Dict[str, Any], pdf_content: Dict[str, Any] = None):
        """Create an image-focused slide with multiple images from the same page"""
        # Add multiple images if present and PDF content is available
        if pdf_content and 'page_images' in slide_data and 'page_number' in slide_data:
            page_number = slide_data['page_number']
            page_images = slide_data['page_images']
            print(f"Adding {len(page_images)} images from page {page_number}")
            
            # Calculate positions for multiple images
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            
            # Position images in a grid layout
            images_per_row = 2
            image_width = Inches(3.5)
            image_height = Inches(2.5)
            
            for i, img_data in enumerate(page_images):
                # Calculate position
                row = i // images_per_row
                col = i % images_per_row
                
                left = Inches(1 + col * 4)  # 1 inch margin + 4 inch spacing
                top = Inches(2 + row * 3)   # 2 inch margin + 3 inch spacing
                
                # Add image to slide
                try:
                    # Convert bytes to PIL Image
                    image = Image.open(io.BytesIO(img_data['data']))
                    
                    # Save temporarily
                    temp_path = f"temp_image_{page_number}_{i}.png"
                    image.save(temp_path)
                    
                    # Add to slide
                    slide.shapes.add_picture(temp_path, left, top, image_width, image_height)
                    
                    # Clean up
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
                        
                    print(f"Added image {i+1} from page {page_number}")
                    
                except Exception as e:
                    print(f"Error adding image {i+1} from page {page_number}: {e}")
        
        # Add image descriptions if present
        if 'image_descriptions' in slide_data and slide_data['image_descriptions']:
            descriptions = slide_data['image_descriptions']
            print(f"Adding {len(descriptions)} image descriptions")
            
            # Add descriptions below images
            for i, description in enumerate(descriptions):
                try:
                    # Position description below corresponding image
                    row = i // 2  # 2 images per row
                    col = i % 2
                    
                    left = Inches(1 + col * 4)
                    top = Inches(4.5 + row * 3)  # Below images
                    width = Inches(3.5)
                    height = Inches(0.5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    p = text_frame.paragraphs[0]
                    p.text = description
                    p.font.size = Pt(10)
                    p.font.italic = True
                    p.font.color.rgb = RGBColor(89, 89, 89)
                    p.alignment = PP_ALIGN.CENTER
                    
                except Exception as e:
                    print(f"Error adding image description {i+1}: {e}")
    
    def _create_mixed_slide(self, slide, slide_data: Dict[str, Any], pdf_content: Dict[str, Any] = None):
        """Create a mixed slide with both text and image"""
        # Add text content first
        self._create_text_slide(slide, slide_data)
        
        # Add image if present and PDF content is available
        if pdf_content and 'image_index' in slide_data:
            image_index = slide_data['image_index']
            print(f"Adding image to mixed slide with index: {image_index}")
            
            # Find the image in PDF content
            for page in pdf_content.get('pages', []):
                if image_index < len(page['images']):
                    image_data = page['images'][image_index]['data']
                    print(f"Found image data for mixed slide, size: {len(image_data)} bytes")
                    self.add_image_from_bytes(slide, image_data, f"image_{image_index}.png")
                    break
    
    def get_presentation_info(self) -> Dict[str, Any]:
        """Get information about the current presentation"""
        return {
            'slide_count': len(self.presentation.slides),
            'slide_width': self.presentation.slide_width,
            'slide_height': self.presentation.slide_height
        }
