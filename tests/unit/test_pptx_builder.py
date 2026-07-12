"""pptx builder: aspect preservation and per-call freshness (the two v1 bugs)."""

import os

from pptx import Presentation
from pptx.util import Inches

from pdfdeck.models import SlideSpec
from pdfdeck.pptx.builder import DeckBuilder, _fit


def test_fit_preserves_aspect_landscape():
    w, h = _fit(1000, 500, Inches(10), Inches(6))
    assert abs((w / h) - 2.0) < 1e-6  # 2:1 preserved
    assert w <= Inches(10) and h <= Inches(6)


def test_fit_preserves_aspect_portrait():
    w, h = _fit(400, 1200, Inches(10), Inches(6))
    assert abs((w / h) - (1 / 3)) < 1e-6
    assert h == Inches(6)  # height-bound


def test_deck_is_fresh_each_call(tmp_path):
    """v1 bug: a cached builder appended to the prior deck. Each build is fresh."""
    builder = DeckBuilder()
    slides = [SlideSpec(index=0, title="T", slide_type="title")]
    a = builder.build(slides, str(tmp_path / "a.pptx"))
    b = builder.build(slides, str(tmp_path / "b.pptx"))
    assert len(list(Presentation(a).slides)) == 1
    assert len(list(Presentation(b).slides)) == 1  # not 2


def test_figure_slide_image_aspect_matches_source(tmp_path):
    # Make a known-aspect image (300x100 => 3:1).
    from PIL import Image
    img = tmp_path / "fig.png"
    Image.new("RGB", (300, 100), (200, 30, 30)).save(img)

    spec = SlideSpec(index=1, title="Figure 1.1", slide_type="figure",
                     figure_ref="p0_r0", image_path=str(img), caption="Figure 1.1 Test.")
    out = DeckBuilder().build([spec], str(tmp_path / "d.pptx"))
    prs = Presentation(out)
    pics = [sh for sh in list(prs.slides)[0].shapes if sh.shape_type == 13]
    assert len(pics) == 1
    ratio = pics[0].width / pics[0].height
    assert abs(ratio - 3.0) < 0.02  # aspect preserved, not stretched to the box
